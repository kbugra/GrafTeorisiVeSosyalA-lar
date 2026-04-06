import sys
import os
sys.stdout.reconfigure(encoding='utf-8')
from pptx import Presentation

files = ['GRAF 1.pptx', 'GRAF 2.pptx', 'Graph Hafta 4.pptx', 'GraphTheory6.pptx', 'Graph 7.pptx']
for f in files:
    fp = os.path.join(r'c:\Users\zxc\Desktop\genel\Üniversite\graf', f)
    if os.path.exists(fp):
        print(f'===== {f} =====')
        prs = Presentation(fp)
        for i, slide in enumerate(prs.slides):
            texts = []
            for shape in slide.shapes:
                if hasattr(shape, 'text_frame'):
                    for para in shape.text_frame.paragraphs:
                        t = para.text.strip()
                        if t:
                            texts.append(t)
            if texts:
                print(f'  Slide {i+1}:')
                for t in texts:
                    print(f'    {t}')
        print()
